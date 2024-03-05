import collections 
import collections.abc
from pptx import Presentation
from pptx.slide import Slide
from datetime import datetime
import os.path
import pandas as pd
from fund_performance import gera_df_performance, tables


template_macro = 'Template.pptx'

# Data de início e fim para o cálculo do retorno
start_date = datetime(2023,12,29)
end_date = datetime(2024,2,29)

def decode_layout(slide: Slide) -> tuple[list, list, list]:
    shape_dict = {}
    for shape in slide.shapes:  
        if shape.is_placeholder:
            phf = shape.placeholder_format
            shape_dict[phf.idx] = (phf.type, str(phf.type))
    pics = [i for i in shape_dict if shape_dict[i][0] == 18]
    bodies = [i for i in shape_dict if shape_dict[i][0] == 2]
    subtitles = [i for i in shape_dict if shape_dict[i][0] == 4]
    return pics, bodies, subtitles

def fill_1_grafico(slide: Slide, flat_lists: dict, i: int, retorno_type: str) -> Slide:
    pics, bodies, subtitles = decode_layout(slide)
    if len(bodies) > 1:
        slide.placeholders[bodies[0]].text = f"Retorno + {retorno_type} + {flat_lists['charts'][i]}"
        slide.placeholders[bodies[1]].text = "gerado em " + datetime.now().strftime("%d-%b-%y") + "\n" + "com dados disponíveis até " + endDate.strftime("%d-%b-%y")
    else:
        slide.placeholders[bodies[0]].text = "carteiras dos concorrentes defasadas em 3 meses"
    
    slide.placeholders[pics[0]].insert_picture(os.path.join(".", "figures", flat_lists["files"][i] + ".png"))
    return slide

if __name__ == "__main__":
    prs = Presentation(template_macro)
    fundos = ["EON", "EVO"]
    save_files = True
    layouts = {"1_grafico": fill_1_grafico, "comps_slide": fill_1_grafico}

    for fundo, slide_top_color in zip(["EVO","EON"], ["gray","blue"]):
        # Gerar o retorno MTD
        df_mtd, df_ytd = gera_df_performance(fundo, save_files)
        
        # Adicionar o slide para o retorno MTD
        slide_mtd = prs.slides.add_slide(prs.slide_layouts.get_by_name("comps_slide_" + slide_top_color))
        slide_mtd.shapes.title.text = f"Retornos MTD"
        fill_1_grafico(slide_mtd, {"charts": ["Performance de Fundos MTD"], "files": [f"{tables[fundo][0]}_retorno_mtd"]}, 0, "MTD")
        
        # Adicionar o slide para o retorno YTD
        slide_ytd = prs.slides.add_slide(prs.slide_layouts.get_by_name("comps_slide_" + slide_top_color))
        slide_ytd.shapes.title.text = f"Retornos YTD"
        fill_1_grafico(slide_ytd, {"charts": ["Performance de Fundos YTD"], "files": [f"{tables[fundo][0]}_retorno_ytd"]}, 0, "YTD")
        
    # Salvar a apresentação
    filename = "Fund_return"
    prs.save(os.path.join("PPT", filename + end_date.strftime("%m-%y") + ".pptx"))
    print("Apresentação salva como " + filename + end_date.strftime("%m-%y") + ".pptx")