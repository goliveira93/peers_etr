"""
Cria report.pptx com a performance de cada fundo dentro da carteira dos nossos FOFs
"""

# pylint: disable=import-error
# pylint: disable=wrong-import-position
# pylint: disable=no-name-in-module
# pylint: disable=invalid-name
import os
import pandas as pd
from datetime import datetime
from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE # type: ignore
from pptx.enum.text import PP_ALIGN  # type: ignore
from summary import get_FOF_price_change  #type: ignore
from settings import britech_data_final

carteira_eon = [
    {"Nome":"Ibiuna STB","Ticker":"27825226000188","Source":"Quantum"},
    {"Nome":"Kapitalo Zeta","Ticker":"12105992000109","Source":"Quantum"},
    {"Nome":"Legacy Alpha","Ticker":"49722651000184","Source":"Quantum", "Proxy ticker":"31666755000153", "Proxy source":"Quantum"},
    {"Nome":"SPX Raptor","Ticker":"26516247000159","Source":"Quantum", "Proxy ticker":"12809201000113", "Proxy source":"Quantum"},
    {"Nome":"SPX Nimitz","Ticker":"12831360000114","Source":"Quantum"},
    {"Nome":"Vinland Macro Plus","Ticker":"30593439000136","Source":"Quantum"},
    {"Nome":"Giant Zarathustra","Ticker":"11052478000181","Source":"Quantum"},
    {"Nome":"Kadima High Vol","Ticker":"14146496000110","Source":"Quantum"},
    {"Nome":"Ibiuna Long short","Ticker":"18391138000124","Source":"Quantum"},
    {"Nome":"RPS","Ticker":"18611600000151","Source":"Quantum"},
    {"Nome":"Absolute Alpha MARB","Ticker":"35618055000144","Source":"Quantum"}
]

carteira_evo = [
    {"Nome":"Encore LB","Ticker":"37487439000109","Source":"Quantum"},
    {"Nome":"3 Ilhas", "Ticker": "51152458000105", "Source":"Quantum"},
    {"Nome":"Atmos","Ticker":"11145320000156","Source":"Quantum"},
    {"Nome":"Dynamo","Ticker":"73232530000139","Source":"Quantum"},
    {"Nome":"Ibiúna LO", "Ticker": "26243348000101", "Source":"Quantum"},
    {"Nome":"Kiron","Ticker":"25213366000170","Source":"Quantum"},
    {"Nome":"Núcleo","Ticker":"37367932000187","Source":"Quantum", "Proxy ticker":"14068366000107", "Proxy source":"Quantum"},
    {"Nome":"Organon","Ticker":"17400251000166","Source":"Quantum"}
]

carteira_outros = [
    {"Nome":"Capitânia Radar 90","Ticker":"23272391000107","Source":"Quantum"},
    {"Nome":"Capitânia Infra 90","Ticker":"27923072000167","Source":"Quantum"},
    {"Nome":"Augme 180","Ticker":"34218678000167","Source":"Quantum"},
    {"Nome":"IMA-B tracker","Ticker":"51514860000184","Source":"Quantum"},
    {"Nome":"IMA-B","Ticker":"IMA-B","Source":"Quantum"},
    {"Nome":"CDI","Ticker":"CDI","Source":"Quantum"}
]


def add_update(my_slide, enddate:str):
    """
    Coloca tarja vermelha com update no slide
    """
    left = Inches(2)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1)

    # Add the diagonal rectangle shape
    shape = my_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )

    # Set the rotation of the shape to create a diagonal orientation
    shape.rotation = -30

    # Set the fill color of the shape
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color

    # Add the text to the shape
    text_frame = shape.text_frame
    text_frame.text = "PRICES UPDATED "+enddate

    # Set the font properties of the text
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    font = run.font
    font.size = Pt(32)
    font.color.rgb = RGBColor(255, 255, 255)  # White color
    return my_slide

def df_to_slide(my_df:pd.DataFrame, my_slide:Slide, my_titulo:str)->Slide:
    """
    Pega um dataframe com dados de performance e cria o slide
    """
    left = Inches(0.5)
    top = Inches(1.79)
    width = Inches(12.47)
    height = Inches(1)
    column_widths=[Inches(2.52),Inches(1.44),Inches(1.44),Inches(1.44),Inches(1.44),Inches(1.44),Inches(2.78)]
    slide.placeholders[12].text=my_titulo                                                      #type:ignore
    slide.placeholders[11].text="Comentários sobre a performance"                           #type:ignore
    my_df["Fundo"]=my_df.index
    my_df["Link"]=["" for _ in my_df.index]
    my_df=my_df[["Fundo","MTD","YTD","12 meses","24 meses","60 meses","Link"]]
    table = my_slide.shapes.add_table(my_df.shape[0] + 1, my_df.shape[1], left, top, width, height).table  #type: ignore
    for i, column in enumerate(my_df.columns):
        table.cell(0, i).text = column

    # Populate the table with DataFrame values
    for row in range(my_df.shape[0]):
        for col in range(my_df.shape[1]):
            table.cell(row + 1, col).text = str(my_df.iloc[row, col])
    for col, width in enumerate(column_widths):
        table.columns[col].width = width

    for rr,row in enumerate(table.rows):
        for cc,cell in enumerate(row.cells):
            # Set text alignment to center
            for paragraph in cell.text_frame.paragraphs:
                if rr>0:
                    paragraph.font.size = Pt(16)
                if 0 < cc <6:
                    paragraph.alignment = PP_ALIGN.CENTER
    return my_slide


if __name__=="__main__":
    end_date=britech_data_final
    end_date=datetime(2024,2,28)
    prs = Presentation(os.path.join(".","Template.pptx"))    # Só precisa abrir o arquivo 1x
    carteira_evo+=[{"Nome":"IBX","Ticker":"IBX","Source":"Quantum"}]
    carteira_eon+=[{"Nome":"IFMM","Ticker":"IFMM BTG PACTUAL","Source":"Quantum"},{"Nome":"CDI","Ticker":"CDI","Source":"Quantum"}]
    for carteira, titulo in zip([carteira_eon, carteira_evo, carteira_outros],["ETRNTY ÉON", "ETRNTY EVO", "OUTROS"]):
        df=get_FOF_price_change(carteira,end_date).sort_index()
        for r in ["IBX","IFMM","CDI"]:
            if r in df.index:
                row_to_move=df.loc[r].copy()
                df=df.drop(r)
                df.loc[r]=row_to_move

        df["MTD"]=df["MTD"].apply(lambda x: str(x)+"%")
        df["YTD"]=df["YTD"].apply(lambda x: str(x)+"%")
        df["12 meses"]=df["12 meses"].apply(lambda x: str(x)+"%")
        df["24 meses"]=df["24 meses"].apply(lambda x: str(x)+"%")
        df["60 meses"]=df["60 meses"].apply(lambda x: str(x)+"%")
        slide = prs.slides.add_slide(prs.slide_layouts.get_by_name("tabela_fundos"))  # You can choose a different layout if needed
        slide=df_to_slide(df,slide,titulo)

    prs.save(os.path.join(".","PPT","report"+".pptx"))
