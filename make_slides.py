import argparse
import collections 
import collections.abc
from pptx import Presentation
from pptx.slide import Slide
from datetime import datetime
import os.path
from libs.delta_etrnty import gera_df
import pandas as pd
from libs.heatmap import make_heatmap
import summary
from conta_reunioes import make_numero_reunioes_fig
from fund_performance import gera_df_performance, tables
#from settings import endDate
import fof

template_macro = 'Template.pptx'

def decode_layout(slide : Slide)-> tuple[list,list,list]:
    shape_dict={}
    for shape in slide.shapes:  #type: ignore
        if shape.is_placeholder:
            phf = shape.placeholder_format
            # print(f"Índice: {phf.idx}, Tipo: {phf.type}, Nome: {shape.name}")
            shape_dict[phf.idx]=(phf.type, str(phf.type))
    pics=[i for i in shape_dict if shape_dict[i][0]==18]
    bodies=[i for i in shape_dict if shape_dict[i][0]==2]
    subtitles=[i for i in shape_dict if shape_dict[i][0]==4]
    return pics,bodies,subtitles

def fill_1_grafico(slide:Slide, flat_lists:dict, i:int, endDate:datetime)->Slide:
    pics,bodies,subtitles=decode_layout(slide)
    if len(bodies)>1:
        slide.placeholders[bodies[0]].text=flat_lists["charts"][i]#type: ignore
        slide.placeholders[bodies[1]].text="gerado em "+datetime.now().strftime("%d-%b-%y")+"\n"+"com dados disponíveis até "+endDate.strftime("%d-%b-%y")  #type: ignore
    else:
        slide.placeholders[bodies[0]].text="carteiras dos concorrentes defasadas em 3 meses"  #type: ignore
    
    slide.placeholders[pics[0]].insert_picture(os.path.join(".","figures",flat_lists["files"][i]+".png"))   #type: ignore
    return slide

def fill_2_graficos(slide: Slide, flat_lists: dict) -> Slide:
    pics, bodies, subtitles = decode_layout(slide)

    slide.placeholders[bodies[2]].text = flat_lists["charts"]["left"]  # type: ignore
    slide.placeholders[bodies[0]].text = flat_lists["charts"]["right"]  # type: ignore
    slide.placeholders[bodies[1]].text = "_contribuição do excesso de performance em relação ao fundo Etrnty"  #type: ignore
    slide.placeholders[bodies[3]].text = " "  # type: ignore

    if flat_lists["files"]["left"] is not None:
        slide.placeholders[pics[0]].insert_picture(os.path.join(".","figures", flat_lists["files"]["left"] + ".png"))  # type: ignore
    if flat_lists["files"]["right"] is not None:
        slide.placeholders[pics[1]].insert_picture(os.path.join(".","figures", flat_lists["files"]["right"] + ".png"))  # type: ignore
    if subtitles:
        slide.placeholders[subtitles[0]].text = f"Etrnty vs {flat_lists['charts']['right']}"                            #type: ignore

    return slide

def fill_performance_comp(slide, args: dict, endDate:datetime) -> Slide:
    pics, bodies, subtitles = decode_layout(slide)
    slide.placeholders[0].text=args["title"]
    if args["files"][0] is not None:
        slide.placeholders[pics[0]].insert_picture(os.path.join(".","figures", args["files"][0] + ".png"))  # type: ignore
    if args["files"][1]is not None:
        slide.placeholders[pics[1]].insert_picture(os.path.join(".","figures", args["files"][1] + ".png"))  # type: ignore
    return slide

def fill_returns(slide: Slide, df_final_ultimo_mes, gestor: str, periodo: str, endDate:datetime) -> Slide:
    # Encontra a linha no dataframe onde a coluna 'Gestor' é igual ao gestor desejado
    filtro = df_final_ultimo_mes['Gestor'] == gestor
    # Localiza os placeholders onde o texto deve ser inserido
    _, bodies, _ = decode_layout(slide)
    if sum(filtro.values==True)==0:
        slide.placeholders[bodies[4]].text = "Não há retornos para o gestor: "+gestor  # type: ignore
        slide.placeholders[bodies[5]].text = "ERRO"  # type: ignore 

    else:
        filtro_ETR = df_final_ultimo_mes['Gestor'] == 'Etrnty'
        
        if periodo == 'MTD':
            coluna_retorno_gestor = 'RETORNO_PEER'
            coluna_retorno_ETR = 'RETORNO_PEER'
        else:  # 'YTD'
            coluna_retorno_gestor = 'RETORNO_COMPOSTO_ACUMULADO_PEER'
            coluna_retorno_ETR = 'RETORNO_COMPOSTO_ACUMULADO_PEER'
        
        retorno_peer = df_final_ultimo_mes.loc[filtro, coluna_retorno_gestor].values[0]
        retorno_ETR = df_final_ultimo_mes.loc[filtro_ETR, coluna_retorno_ETR].values[0]

        # Formata o valor de retorno como uma string de porcentagem
        texto_retorno_peer = f"{retorno_peer:.1%}"
        texto_retorno_ETR = f"{retorno_ETR:.1%}"

        # Insere o texto no placeholder correspondente
        slide.placeholders[bodies[4]].text = texto_retorno_ETR  # type: ignore
        slide.placeholders[bodies[5]].text = texto_retorno_peer  # type: ignore    

    return slide


def cria_slides(endDate:datetime):
    slide = prs.slides.add_slide(prs.slide_layouts.get_by_name("cover_fundos"))
    slide = prs.slides.add_slide(prs.slide_layouts.get_by_name("1_grafico_cinza"))
    slide.shapes.title.text = "SELEÇÃO DE GESTORES"
    fill_1_grafico(slide,{"charts":["Número de interações com gestores"],"files":["reunioes_mes"]},0,endDate)

    for fund, tipo, prefix, slide_top_color, title, template in zip(["EON","EVO"],["Multimercado","Ações"],["FIM","FIA"],["blue","gray"],["ETRNTY EON","ETRNTY EVO"],["1_grafico_azul","1_grafico_cinza"]):
        slide = prs.slides.add_slide(prs.slide_layouts.get_by_name(template))
        slide.shapes.title.text = title
        fill_1_grafico(slide,{"charts":["Alterações no portfólio​"],"files":[fund+"_weight_cng"]},0,endDate)
        slide = prs.slides.add_slide(prs.slide_layouts.get_by_name(template))
        slide.shapes.title.text = title
        fill_1_grafico(slide,{"charts":["Principais performances no mês​​"],"files":[fund+"_price_cngs"]},0,endDate)
        slide = prs.slides.add_slide(prs.slide_layouts.get_by_name(template))
        slide.shapes.title.text = title
        fill_1_grafico(slide,{"charts":["Contribuições para o retorno do mês​​"],"files":[fund+"_contribution"]},0,endDate)
        
        #Pega .png dos gráficos de barra (performance absoluta) e coloca no pptx
        slide = prs.slides.add_slide(prs.slide_layouts.get_by_name("performance_comp_"+slide_top_color))
        fill_performance_comp(slide,{"files":[fund.lower()+"_MTD",fund.lower()+"_YTD"],"title":"ETRNTY "+fund},endDate)

        #Gera heatimap com posições dos concorrentes
        slide = prs.slides.add_slide(prs.slide_layouts.get_by_name("comps_slide_"+slide_top_color))
        slide.shapes.title.text = "PEERS - "+fund
        fill_1_grafico(slide,{"charts":["carteira pares"],"files":["heatmap_"+fund]},0,endDate)

        #gera comparações
        for period in ["MTD","YTD"]:
            #slide com as peformances dos fundos dentro dos FoFs da concorrência
            slide = prs.slides.add_slide(prs.slide_layouts.get_by_name("comps_slide_" + slide_top_color))
            slide.shapes.title.text = "Retornos "+period
            fill_1_grafico(slide, {"charts": ["Performance de Fundos "+period], "files": [f"{tables[fund][0]}_retorno_"+period.lower()]}, 0,endDate)

            try:
                df_final_ultimo_mes=gera_df(fund,period)
            except Exception as e:
                print("Erro gerando gráficos para: "+fund)
                break

            for gestor in gestores:
                if gestor != "Etrnty":
                    # Define os nomes dos arquivos para MTD e YTD
                    arquivo_gestor = prefix+"_"+gestor+"_"+period
                    arquivo_etrnty = prefix+"_Etrnty"+"_"+period
                    
                    # Adiciona um novo slide para MTD
                    slide = prs.slides.add_slide(prs.slide_layouts.get_by_name("2_graficos"))
                    # Preencher o slide com dados de MTD
                    slide = fill_returns(slide, df_final_ultimo_mes, gestor, period,endDate)
                    # Configurar títulos, gráficos, e outras informações para MTD e YTD
                    slide.shapes.title.text = tipo+" - "+period                                #type: ignore
                    
                    # Preencher os slides com os dados
                    flat_lists = {
                        "charts": {"left": "Etrnty", "right": gestor},
                        "files": {"left": arquivo_etrnty, "right": arquivo_gestor}
                    }
                    
                    # Preencher o slide com layouts específicos
                    layouts["2_graficos"](slide, flat_lists)


if __name__=="__main__":
    endDate=datetime.strptime("30092024","%d%m%Y")
    startDate=datetime.strptime("30082024","%d%m%Y")
    YTD_base_db=datetime.strptime("31012024","%d%m%Y")
    YTD_date=datetime(2023,12,29)
    filename=""
    layouts = {"1_grafico": fill_1_grafico, "2_graficos": fill_2_graficos}
    gestores = ["Brain", "Consenso", "Etrnty", "G5", "JBFO", "Mandatto", "Portofino", "Pragma", "Warren", "We Capital", "Wright", "XPA"]

    parser = argparse.ArgumentParser(description='make_slides')
    # Adicione argumentos usando o método add_argument
    parser.add_argument('--parcial', choices=["resumo","analitico","fof","apenas_slides"], default=None,help="Se especificado, roda apenas parte do relatório" )
    args = parser.parse_args()
    
    if args.parcial=="apenas_slides":
        filename="slides_"
    else:
        if args.parcial is None or args.parcial=="resumo":
            filename="resumo_"
            figs=[]
            #faz grafico com numero de reunioes
            figs.append(make_numero_reunioes_fig(endDate))
            #faz graficos de barra com performance absoluta YTD, MTD para eon e evo
            try:
                figs=figs+summary.make_summary_figs(endDate,gestores,YTD_date)
            except ValueError as e:
                print(e)
                pass
            if args.parcial is None:
                for f in figs:
                    f.show()  
        if args.parcial is None or args.parcial=="fof":
            figs=[]
            figs=figs+fof.performance_attrib_fof("EON",startDate.date(),endDate.date())  #type: ignore
            figs=figs+fof.performance_attrib_fof("EVO",startDate.date(),endDate.date())  #type: ignore
            if args.parcial is not None:
                for f in figs:
                    f.show()
            

        if args.parcial is None or args.parcial=="analitico":
            filename="analitico_"
            save_files = True

            for fund, tipo, prefix, slide_top_color in zip(["EON","EVO"],["Multimercado","Ações"],["FIM","FIA"],["blue","gray"]):
                #df com as peformances dos fundos dentro dos FoFs da concorrência
                df_mtd, df_ytd = gera_df_performance(fund,YTD_base_db,endDate, save_files)
                #Gera heatimap com posições dos concorrentes
                df=gera_df(fund,"MTD",False)
                df=df[df["Gestor"].isin(gestores)]
                make_heatmap(fund,df)
                
           
    # Salvar a apresentação
    prs = Presentation(template_macro)
    cria_slides(endDate)
    if args.parcial is None:
        filename = "Comparacao_"
    prs.save(os.path.join("PPT", filename + endDate.strftime("%m-%y") + ".pptx"))
    print("Apresentação salva como " + filename + endDate.strftime("%m-%y") + ".pptx")
    print("Lembre de rodar o make_report.py agora.")

